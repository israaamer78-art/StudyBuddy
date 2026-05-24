"""Parse notes from PDF, DOCX, or plain text files.

For PDFs, we render each page as an image and send it to the selected vision API
to extract ALL text including labels embedded inside diagrams and images —
which regular PDF text extraction misses entirely.
"""
import hashlib
import io
from pathlib import Path

import ai_provider
import cache_store
from model_config import current_model

CACHE_DIR = cache_store.CACHE_DIR
NOTES_CACHE_VERSION = 1
PPTX_IMAGE_VISION_VERSION = 1
MAX_PPTX_VISION_IMAGES = 20

PDF_VISION_PROMPT = (
    "You are extracting content from a medical school lecture slide. "
    "First, look at the slide and decide what type it is, then extract accordingly.\n\n"
    "**SLIDE TYPE A: Text-as-image** (the slide is mostly a screenshot of text, a table, or styled text boxes — like a screenshot from a textbook or website)\n"
    "→ Transcribe EVERY WORD exactly as it appears. Do not paraphrase. Do not skip rows. "
    "If there's a table, format it as a proper markdown table with | separators, preserving every cell.\n\n"
    "**SLIDE TYPE B: Text content with a decorative or supporting diagram** (the slide has bullet points, colored boxes, headers, AND a diagram — but the text content carries the learning point)\n"
    "→ Transcribe EVERY word of the text content (bullets, headers, colored boxes, captions). "
    "For the diagram, write a brief 1-2 sentence description of what it shows. "
    "Do NOT list every label on the diagram unless a label is clearly central (e.g. highlighted, called out by an arrow with emphasis).\n\n"
    "**SLIDE TYPE C: Dense anatomical diagram with many labels** (the main content IS the labeled diagram — like an atlas figure with 20+ labels)\n"
    "→ Describe what views/structures the diagram shows in 1-2 sentences. "
    "List ONLY the labels that appear central to the learning objective (e.g. larger text, the slide's title topic, items the lecturer would emphasize). "
    "Do NOT list every minor label — a study guide doesn't need 'lateral mass' AND 'inferolateral angle' AND 'S2' AND 'S3' AND 'S4' AND 'S5' all itemized. "
    "Pick out maybe 5-10 of the most clinically important structures shown.\n\n"
    "**SLIDE TYPE D: Title slide or section divider** (just a topic name on a mostly empty slide)\n"
    "→ Write just the title. One line.\n\n"
    "**REGARDLESS OF TYPE:**\n"
    "- Extract the slide title and any headers exactly as shown.\n"
    "- Transcribe any handwritten or typed annotations the lecturer added.\n"
    "- If there's a star (★) or highlight or 'IMPORTANT' marker, note it.\n"
    "- Do not add commentary like 'this slide shows' or 'we can see that'.\n"
    "- Do not invent content not visible on the slide.\n\n"
    "Begin extraction now:"
)

PPTX_IMAGE_PROMPT = (
    "Extract the educational content from this PowerPoint slide image or embedded diagram. "
    "Focus on visible labels, table text, relationships, arrows, highlighted structures, and clinically relevant facts. "
    "Do not invent details that are not visible. Be concise but complete."
)


def _cached_file_parse(path: Path, kind: str, payload: dict, parser, progress=None) -> str:
    key = cache_store.make_key(kind, {
        "version": NOTES_CACHE_VERSION,
        "file_hash": cache_store.hash_file(path),
        "suffix": path.suffix.lower(),
        **payload,
    })
    cached = cache_store.get(CACHE_DIR, kind, key)
    if cached is not None:
        if progress:
            progress({"stage": "Using cached notes", "section_title": path.name})
        return cached
    value = parser()
    cache_store.set(CACHE_DIR, kind, key, value)
    return value


def parse_pdf_text_only(path: str) -> str:
    """Fast fallback: only extract text that's stored as text in the PDF."""
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def parse_pdf_with_vision(path: str, max_pages: int = 50, progress=None) -> str:
    """Render each PDF page as an image, send to vision for OCR + text extraction."""
    from pypdf import PdfReader
    from pdf2image import convert_from_path
    reader = PdfReader(path)
    page_count = len(reader.pages)

    if page_count > max_pages:
        print(f"⚠️  PDF has {page_count} pages, processing first {max_pages} via vision...")
        page_count = max_pages

    print(f"📸 Rendering {page_count} page(s) and extracting via vision...")
    if progress:
        progress({"stage": "Loading notes from PDF", "item_label": "page", "current": 0, "total": page_count})

    extracted = []
    for page_num in range(1, page_count + 1):
        try:
            if progress:
                progress({
                    "stage": "Loading notes from PDF",
                    "item_label": "page",
                    "current": page_num,
                    "total": page_count,
                    "section_title": f"Page {page_num}",
                })
            images = convert_from_path(path, dpi=150, first_page=page_num, last_page=page_num)
            if not images:
                page_text = reader.pages[page_num - 1].extract_text() or ""
                extracted.append(f"--- Slide {page_num} ---\n{page_text}")
                continue

            img = images[0]
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            text = ai_provider.call_vision_text(
                buf.getvalue(),
                "image/png",
                PDF_VISION_PROMPT,
                max_tokens=2048,
            )
            extracted.append(f"--- Slide {page_num} ---\n{text}")
            print(f"   [{page_num}/{page_count}] ✓")
            if progress:
                progress({
                    "stage": "Loading notes from PDF",
                    "item_label": "page",
                    "current": page_num,
                    "total": page_count,
                    "completed": page_num,
                    "section_title": f"Page {page_num}",
                })

        except Exception as e:
            print(f"   [{page_num}/{page_count}] ⚠️  Vision failed, using text fallback: {e}")
            try:
                page_text = reader.pages[page_num - 1].extract_text() or ""
                extracted.append(f"--- Slide {page_num} ---\n{page_text}")
            except Exception:
                extracted.append(f"--- Slide {page_num} ---\n(extraction failed)")

    return "\n\n".join(extracted)


def parse_pdf(path: str, use_vision: bool = True, progress=None) -> str:
    if use_vision:
        try:
            return parse_pdf_with_vision(path, progress=progress)
        except ImportError as e:
            print(f"⚠️  Vision dependency missing ({e}). Falling back to plain text extraction.")
            print("   To enable vision: pip install pdf2image, plus install poppler (brew install poppler on Mac)")
            if progress:
                progress({"stage": "Loading notes from PDF text"})
            return parse_pdf_text_only(path)
        except Exception as e:
            print(f"⚠️  Vision PDF parsing failed: {e}. Falling back to plain text.")
            if progress:
                progress({"stage": "Loading notes from PDF text"})
            return parse_pdf_text_only(path)
    else:
        if progress:
            progress({"stage": "Loading notes from PDF text"})
        return parse_pdf_text_only(path)


def parse_docx(path: str, progress=None) -> str:
    if progress:
        progress({"stage": "Loading DOCX notes", "section_title": Path(path).name})
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def parse_txt(path: str, progress=None) -> str:
    if progress:
        progress({"stage": "Loading text notes", "section_title": Path(path).name})
    return Path(path).read_text(encoding="utf-8")


def analyze_pptx_image_with_vision(image_bytes: bytes, content_type: str | None = None) -> str:
    media_type = content_type or "image/png"
    return ai_provider.call_vision_text(image_bytes, media_type, PPTX_IMAGE_PROMPT, max_tokens=1024)


def cached_pptx_image_analysis(image_bytes: bytes, content_type: str | None = None, progress=None, label: str = "") -> str:
    image_hash = hashlib.sha256(image_bytes).hexdigest()
    key = cache_store.make_key("pptx_image_vision", {
        "version": PPTX_IMAGE_VISION_VERSION,
        "model": current_model(),
        "image_hash": image_hash,
        "content_type": content_type or "image/png",
    })
    cached = cache_store.get(CACHE_DIR, "pptx_image_vision", key)
    if cached is not None:
        if progress:
            progress({"stage": "Using cached PPTX image analysis", "section_title": label})
        return cached
    if progress:
        progress({"stage": "Analyzing PPTX embedded image", "section_title": label})
    value = analyze_pptx_image_with_vision(image_bytes, content_type)
    cache_store.set(CACHE_DIR, "pptx_image_vision", key, value)
    return value


def iter_picture_shapes(shapes):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_picture_shapes(shape.shapes)


def parse_pptx(path: str, progress=None, analyze_images: bool = True, max_vision_images: int = MAX_PPTX_VISION_IMAGES) -> str:
    """Extract text from PPTX. For PowerPoint, we extract slide text directly
    rather than rendering to images — python-pptx pulls text from text boxes
    AND from text frames inside shapes/diagrams, which catches most labels."""
    from pptx import Presentation
    prs = Presentation(path)
    slides_text = []
    slide_count = len(prs.slides)
    analyzed_images = 0
    if progress:
        progress({"stage": "Loading PPTX notes", "item_label": "slide", "current": 0, "total": slide_count})
    for i, slide in enumerate(prs.slides, 1):
        if progress:
            progress({
                "stage": "Loading PPTX notes",
                "item_label": "slide",
                "current": i,
                "total": slide_count,
                "section_title": f"Slide {i}",
            })
        parts = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
            elif hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        if analyze_images and analyzed_images < max_vision_images:
            for image_index, shape in enumerate(iter_picture_shapes(slide.shapes), 1):
                if analyzed_images >= max_vision_images:
                    parts.append(f"[Skipped image analysis: reached {max_vision_images} embedded images for this deck]")
                    break
                try:
                    image = shape.image
                    analysis = cached_pptx_image_analysis(
                        image.blob,
                        image.content_type,
                        progress=progress,
                        label=f"Slide {i} image {image_index}",
                    )
                    if analysis:
                        parts.append(f"[Image {image_index} analysis: {analysis}]")
                    analyzed_images += 1
                except Exception as e:
                    parts.append(f"[Image {image_index} analysis failed: {e}]")
        # Notes pane
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(f"[Speaker notes: {notes_text}]")
        slides_text.append("\n".join(parts))
        if progress:
            progress({
                "stage": "Loading PPTX notes",
                "item_label": "slide",
                "current": i,
                "total": slide_count,
                "completed": i,
                "section_title": f"Slide {i}",
            })
    return "\n\n".join(slides_text)


def load_notes(path: str, use_vision: bool = True, progress=None) -> str:
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"Notes file not found: {path}")
    ext = p.suffix.lower()
    print(f"📓 Loading notes from {p.name}...")
    if progress:
        progress({"stage": "Loading notes", "section_title": p.name})
    if ext == ".pdf":
        return _cached_file_parse(
            p,
            "notes_pdf",
            {"use_vision": use_vision, "model": current_model()},
            lambda: parse_pdf(str(p), use_vision=use_vision, progress=progress),
            progress=progress,
        )
    elif ext == ".pptx":
        return _cached_file_parse(
            p,
            "notes_pptx",
            {},
            lambda: parse_pptx(str(p), progress=progress),
            progress=progress,
        )
    elif ext == ".docx":
        return _cached_file_parse(
            p,
            "notes_docx",
            {},
            lambda: parse_docx(str(p), progress=progress),
            progress=progress,
        )
    elif ext in (".txt", ".md"):
        return _cached_file_parse(
            p,
            "notes_text",
            {},
            lambda: parse_txt(str(p), progress=progress),
            progress=progress,
        )
    else:
        raise ValueError(f"Unsupported notes file type: {ext}")
