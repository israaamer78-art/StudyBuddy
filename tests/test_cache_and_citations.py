import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

import generator
import ingest
import notes_parser


def make_tiny_png_bytes():
    import base64
    return base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
    )


class CacheAndCitationTests(unittest.TestCase):
    def test_load_notes_reuses_cached_file_extraction_by_file_hash(self):
        with tempfile.TemporaryDirectory() as tmp:
            notes_parser.CACHE_DIR = Path(tmp)
            path = Path(tmp) / "notes.txt"
            path.write_text("source notes", encoding="utf-8")

            with patch("notes_parser.parse_txt", side_effect=["parsed once", "parsed twice"]) as parse_txt:
                self.assertEqual(notes_parser.load_notes(str(path)), "parsed once")
                self.assertEqual(notes_parser.load_notes(str(path)), "parsed once")
                self.assertEqual(parse_txt.call_count, 1)

    def test_grouping_and_section_generation_cache_uses_model_and_source_hash(self):
        with tempfile.TemporaryDirectory() as tmp:
            generator.CACHE_DIR = Path(tmp)
            slides = [
                {"slide_num": 1, "content": "Pelvis overview"},
                {"slide_num": 2, "content": "Sacral plexus"},
            ]
            grouped = [{"title": "Pelvis", "slide_numbers": [1, 2]}]
            generated = {
                "title": "Pelvis",
                "reading": "reading",
                "key_terms": [],
                "matching": [],
                "diagram": None,
            }

            with patch("generator._call_claude_json", side_effect=[grouped, generated]) as call:
                first_groups = generator.group_slides_into_sections(slides)
                second_groups = generator.group_slides_into_sections(slides)
                self.assertEqual(first_groups, second_groups)

                first_section = generator.generate_section("slides", "transcript", 1, "Pelvis")
                second_section = generator.generate_section("slides", "transcript", 1, "Pelvis")
                self.assertEqual(first_section["reading"], second_section["reading"])
                self.assertEqual(call.call_count, 2)

    def test_sections_include_source_citations_for_slides_and_transcript(self):
        group = {
            "title": "Pelvis",
            "slide_numbers": [3, 4],
            "slides_content": "--- Slide 3 ---\nPelvis\n\n--- Slide 4 ---\nSacral plexus",
            "transcript_excerpt": "The professor explains the sacral plexus in context.",
            "transcript_index": 2,
        }

        citations = generator.build_source_citations(group)

        self.assertEqual(citations["slides"][0]["slide_number"], 3)
        self.assertEqual(citations["slides"][1]["slide_number"], 4)
        self.assertEqual(citations["transcript"]["chunk_index"], 2)
        self.assertIn("sacral plexus", citations["transcript"]["excerpt"].lower())

    def test_parse_pptx_appends_cached_embedded_image_analysis(self):
        from pptx import Presentation
        from pptx.util import Inches

        with tempfile.TemporaryDirectory() as tmp:
            notes_parser.CACHE_DIR = Path(tmp) / "cache"
            image_path = Path(tmp) / "image.png"
            image_path.write_bytes(make_tiny_png_bytes())

            pptx_path = Path(tmp) / "slides.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            box.text = "Visible slide text"
            slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), width=Inches(1))
            prs.save(pptx_path)

            with patch("notes_parser.analyze_pptx_image_with_vision", return_value="Image shows a labeled pelvis") as analyze:
                first = notes_parser.parse_pptx(str(pptx_path))
                second = notes_parser.parse_pptx(str(pptx_path))

            self.assertIn("Visible slide text", first)
            self.assertIn("[Image 1 analysis: Image shows a labeled pelvis]", first)
            self.assertEqual(first, second)
            self.assertEqual(analyze.call_count, 1)

    def test_pdf_lecture_source_uses_vision_parser(self):
        with tempfile.TemporaryDirectory() as tmp:
            ingest.CACHE_DIR = Path(tmp) / "cache"
            pdf_path = Path(tmp) / "slides.pdf"
            pdf_path.write_bytes(b"pdf")
            progress_events = []

            with patch("notes_parser.parse_pdf", return_value="--- Slide 1 ---\nDiagram labels") as parse_pdf:
                text = ingest.get_transcript(str(pdf_path), progress=progress_events.append)

            self.assertEqual(text, "--- Slide 1 ---\nDiagram labels")
            parse_pdf.assert_called_once_with(str(pdf_path), use_vision=True, progress=progress_events.append)


if __name__ == "__main__":
    unittest.main()
